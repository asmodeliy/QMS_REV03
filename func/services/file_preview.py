
from pathlib import Path 
from openpyxl import load_workbook 
from docx import Document 
from docx .enum .text import WD_ALIGN_PARAGRAPH 
try:
    from pypdf import PdfReader
except ImportError:
    from PyPDF2 import PdfReader 
from pptx import Presentation 
import base64 
import io 
from PIL import Image 


def get_run_style_css (run )->str :
    styles =[]

    if run .bold :
        styles .append ("font-weight: bold;")
    if run .italic :
        styles .append ("font-style: italic;")
    if run .underline :
        styles .append ("text-decoration: underline;")

    if run .font .color and run .font .color .rgb :
        color =str (run .font .color .rgb )
        if color and color !="00000000":
            styles .append (f"color: #{color [-6 :]};")

    if run .font .size :
        size_pt =run .font .size .pt 
        styles .append (f"font-size: {size_pt }pt;")

    return "".join (styles )


def extract_images_from_word_part (doc_part )->dict :
    images ={}
    try :
        for rel_id ,rel in doc_part .rels .items ():
            if "image"in rel .reltype :
                try :
                    image_part =rel .target_part 
                    image_bytes =image_part .blob 
                    base64_image =base64 .b64encode (image_bytes ).decode ('utf-8')

                    ext =image_part .partname .lower ().split ('.')[-1 ]
                    mime_types ={
                    'jpg':'image/jpeg','jpeg':'image/jpeg','png':'image/png',
                    'gif':'image/gif','bmp':'image/bmp','webp':'image/webp',
                    }
                    mime_type =mime_types .get (ext ,'image/jpeg')
                    images [rel_id ]=f"data:{mime_type };base64,{base64_image }"
                except :
                    pass 
    except :
        pass 
    return images 


def preview_word (file_path :Path )->str :
    """Word 문서 미리보기"""
    try :
        doc =Document (file_path )
        html_parts =[]
        images =extract_images_from_word_part (doc .part )

        for para in doc .paragraphs [:200 ]:
            para_html =""
            style_name =para .style .name if para .style else "Normal"

            align_css =""
            if para .alignment ==WD_ALIGN_PARAGRAPH .CENTER :
                align_css ="text-align: center;"
            elif para .alignment ==WD_ALIGN_PARAGRAPH .RIGHT :
                align_css ="text-align: right;"
            elif para .alignment ==WD_ALIGN_PARAGRAPH .JUSTIFY :
                align_css ="text-align: justify;"

            if style_name .startswith ('Heading'):
                level ='2'
                try :
                    if len (style_name )>7 :
                        level =style_name .split ()[-1 ]if style_name .split ()[-1 ].isdigit ()else '2'
                except :
                    pass 

                para_html =f'<h{level } style="{align_css }">'
                for run in para .runs :
                    run_css =get_run_style_css (run )
                    if run_css :
                        para_html +=f'<span style="{run_css }">{run .text }</span>'
                    else :
                        para_html +=run .text 

                    for drawing in run ._element .findall ('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing'):
                        try :
                            blip =drawing .find ('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                            if blip is not None :
                                embed =blip .get ('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                                if embed and embed in images :
                                    para_html +=f'<br><img src="{images [embed ]}" style="max-width: 100%; height: auto; margin: 10px 0; border-radius: 4px;">'
                        except :
                            pass 
                para_html +=f'</h{level }>'

            elif para .text .strip ():
                para_html =f'<p style="margin: 10px 0; line-height: 1.6; {align_css }">'
                for run in para .runs :
                    run_css =get_run_style_css (run )
                    if run_css :
                        para_html +=f'<span style="{run_css }">{run .text }</span>'
                    else :
                        para_html +=run .text 

                    for drawing in run ._element .findall ('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing'):
                        try :
                            blip =drawing .find ('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                            if blip is not None :
                                embed =blip .get ('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                                if embed and embed in images :
                                    para_html +=f'<br><img src="{images [embed ]}" style="max-width: 100%; height: auto; margin: 10px 0; border-radius: 4px;">'
                        except :
                            pass 
                para_html +='</p>'
            else :
                for run in para .runs :
                    for drawing in run ._element .findall ('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing'):
                        try :
                            blip =drawing .find ('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                            if blip is not None :
                                embed =blip .get ('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                                if embed and embed in images :
                                    para_html +=f'<p style="text-align: center;"><img src="{images [embed ]}" style="max-width: 100%; height: auto; margin: 10px 0; border-radius: 4px;"></p>'
                        except :
                            pass 

            if para_html :
                html_parts .append (para_html )

        for table_idx ,table in enumerate (doc .tables [:20 ]):
            html_parts .append ('<table style="border-collapse: collapse; width: 100%; margin: 20px 0; border: 1px solid #d1d5db;">')
            for row_idx ,row in enumerate (table .rows ):
                html_parts .append ('<tr>')
                for cell in row .cells :
                    bg_color ="background-color: #f3f4f6;"if row_idx ==0 else ""
                    html_parts .append (f'<td style="border: 1px solid #d1d5db; padding: 12px; {bg_color }">')
                    for para in cell .paragraphs :
                        if para .text .strip ():
                            html_parts .append (f'<p style="margin: 0;">{para .text }</p>')
                        for run in para .runs :
                            for drawing in run ._element .findall ('.//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}drawing'):
                                try :
                                    blip =drawing .find ('.//{http://schemas.openxmlformats.org/drawingml/2006/main}blip')
                                    if blip is not None :
                                        embed =blip .get ('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                                        if embed and embed in images :
                                            html_parts .append (f'<img src="{images [embed ]}" style="max-width: 100%; height: auto; margin: 10px 0; border-radius: 4px;">')
                                except :
                                    pass 
                    html_parts .append ('</td>')
                html_parts .append ('</tr>')
            html_parts .append ('</table>')

        if len (doc .paragraphs )>200 :
            html_parts .append (f'<p style="color:#6b7280; margin-top: 20px;">... 나머지 {len (doc .paragraphs )-200 }개 문단 생략</p>')

        return '\n'.join (html_parts )if html_parts else "<p>내용이 없습니다.</p>"
    except Exception as e :
        return f"<p>Word 파일을 읽을 수 없습니다: {str (e )}</p>"


def preview_excel (file_path :Path )->str :
    try :
        from openpyxl import load_workbook 

        wb =load_workbook (file_path ,data_only =False )
        total_sheets =len (wb .sheetnames )
        max_sheets =min (total_sheets ,5 )

        html_parts =[]
        html_parts .append (f'<div style="color: #6b7280; font-size: 12px; margin-bottom: 15px;">총 {total_sheets }개 시트 (처음 {max_sheets }개 시트 표시)</div>')

        for sheet_name in wb .sheetnames [:max_sheets ]:
            ws =wb [sheet_name ]

            html_parts .append (f'<h3 style="margin-top: 30px; color: #1f2937; font-weight: 600;">📊 {sheet_name }</h3>')
            html_parts .append ('<div style="overflow-x: auto; max-width: 100%;"><table style="border-collapse: collapse; font-size: 13px; min-width: 600px;">')

            max_row =min (ws .max_row ,100 )
            max_col =min (ws .max_column ,50 )

            for row_idx ,row in enumerate (ws .iter_rows (min_row =1 ,max_row =max_row ,min_col =1 ,max_col =max_col ),start =1 ):
                html_parts .append ('<tr>')

                for col_idx ,cell in enumerate (row ,start =1 ):
                    is_merged =False 
                    skip_cell =False 
                    rowspan =1 
                    colspan =1 

                    for merged_range in ws .merged_cells .ranges :
                        if cell .coordinate in merged_range :
                            is_merged =True 
                            if cell .coordinate !=merged_range .start_cell .coordinate :
                                skip_cell =True 
                                break 
                            else :
                                rowspan =merged_range .max_row -merged_range .min_row +1 
                                colspan =merged_range .max_col -merged_range .min_col +1 
                            break 

                    if skip_cell :
                        continue 

                    value =cell .value 
                    if value is None :
                        value =''
                    else :
                        value =str (value )

                    bg_hex ='#ffffff'
                    try :
                        if cell .fill and cell .fill .start_color :
                            color =cell .fill .start_color 
                            if hasattr (color ,'rgb')and color .rgb :
                                rgb =str (color .rgb )
                                if rgb and rgb !='00000000':
                                    if len (rgb )==8 :
                                        bg_hex =f'#{rgb [2 :]}'
                                    elif len (rgb )==6 :
                                        bg_hex =f'#{rgb }'
                    except :
                        pass 

                    font_hex ='#000000'
                    font_weight ='normal'
                    font_style ='normal'
                    font_size ='13px'

                    try :
                        if cell .font :
                            if cell .font .color and hasattr (cell .font .color ,'rgb')and cell .font .color .rgb :
                                rgb =str (cell .font .color .rgb )
                                if rgb and rgb !='00000000':
                                    if len (rgb )==8 :
                                        font_hex =f'#{rgb [2 :]}'
                                    elif len (rgb )==6 :
                                        font_hex =f'#{rgb }'

                            if cell .font .bold :
                                font_weight ='bold'

                            if cell .font .italic :
                                font_style ='italic'

                            if cell .font .size :
                                font_size =f'{int (cell .font .size )}px'
                    except :
                        pass 

                    text_align ='left'
                    try :
                        if cell .alignment and cell .alignment .horizontal :
                            align_map ={'center':'center','right':'right','left':'left'}
                            text_align =align_map .get (cell .alignment .horizontal ,'left')
                    except :
                        pass 

                    border_style ='border: 1px solid #d1d5db;'

                    if row_idx ==1 :
                        if bg_hex =='#ffffff':
                            bg_hex ='#f3f4f6'
                        if font_weight =='normal':
                            font_weight ='600'

                    merge_attrs =''
                    if is_merged :
                        if rowspan >1 :
                            merge_attrs +=f' rowspan="{rowspan }"'
                        if colspan >1 :
                            merge_attrs +=f' colspan="{colspan }"'

                    html_parts .append (f'<td{merge_attrs } style="{border_style } background-color: {bg_hex }; color: {font_hex }; font-weight: {font_weight }; font-style: {font_style }; font-size: {font_size }; text-align: {text_align }; padding: 8px; white-space: nowrap;">{value }</td>')

                html_parts .append ('</tr>')

            html_parts .append ('</table></div>')

            if ws .max_row >100 :
                html_parts .append (f'<p style="color:#6b7280;font-size:12px; margin-top: 10px;">... 나머지 {ws .max_row -100 }행 생략</p>')

        if total_sheets >max_sheets :
            html_parts .append (f'<p style="color:#6b7280; text-align: center; margin-top: 30px; padding: 20px; background: linear-gradient(135deg, #10b981 0%, #059669 100%); border-radius: 12px; color: white; font-weight: 500; font-size: 14px;">📚 전체 {total_sheets }개 시트 중 {max_sheets }개만 표시됩니다</p>')

        return '\n'.join (html_parts )

    except Exception as e :
        return f"<p style='color: #dc2626;'>Excel 파일을 읽을 수 없습니다: {str (e )}</p>"


def preview_pdf (file_path :Path )->str :
    try :
        import fitz 

        doc =fitz .open (file_path )
        total_pages =len (doc )
        max_pages =min (total_pages ,10 )

        html_parts =[]
        html_parts .append (f'<div style="color: #6b7280; font-size: 12px; margin-bottom: 15px;">총 {total_pages }페이지 (처음 {max_pages }페이지 표시)</div>')

        for page_idx in range (max_pages ):
            try :
                page =doc [page_idx ]

                pix =page .get_pixmap (matrix =fitz .Matrix (2 ,2 ),alpha =False )
                img_data =pix .tobytes ("png")
                base64_image =base64 .b64encode (img_data ).decode ('utf-8')

                html_parts .append (f'<div class="pdf-page" style="margin: 20px 0; text-align: center; border: 1px solid #e5e7eb; padding: 10px; border-radius: 8px;">')
                html_parts .append (f'<div class="page-number" style="font-weight: 600; color: #2e5bff; margin-bottom: 10px;">📄 페이지 {page_idx +1 }</div>')
                html_parts .append (f'<img src="data:image/png;base64,{base64_image }" style="max-width: 100%; height: auto; border-radius: 4px;">')
                html_parts .append ('</div>')
            except Exception as e :
                html_parts .append (f'<div class="pdf-page" style="margin: 20px 0; padding: 15px; border: 1px solid #fee2e2; border-radius: 8px; background: #fef2f2;"><p style="color: #dc2626;">페이지 {page_idx +1 } 렌더링 오류: {str (e )}</p></div>')

        if total_pages >max_pages :
            html_parts .append (f'<p style="color:#6b7280; text-align: center; margin-top: 20px; padding: 15px; background: #f3f4f6; border-radius: 8px;">... 나머지 {total_pages -max_pages }페이지 생략</p>')

        doc .close ()

        return '\n'.join (html_parts )if html_parts else "<p>PDF 내용이 없습니다.</p>"
    except ImportError :
        return "<p style='color: #dc2626;'>PyMuPDF가 설치되지 않았습니다. 'pip install PyMuPDF' 명령어로 설치하세요.</p>"
    except Exception as e :
        return f"<p style='color: #dc2626;'>PDF 파일을 읽을 수 없습니다: {str (e )}</p>"


def preview_pptx (file_path :Path )->str :
    try :
        import tempfile 
        import os 
        import sys 
        import subprocess 
        import glob 
        import time as _ltime 
        from pathlib import Path as PathlibPath 

        prs =Presentation (file_path )
        total_slides =len (prs .slides )
        max_slides =min (total_slides ,100 )


        if sys .platform =='win32':
            powerpoint =None 
            ppt =None 
            temp_dir =tempfile .mkdtemp ()
            html_parts =[]

            try :
                try :
                    import win32com .client 
                    import pythoncom 

                    pythoncom .CoInitialize ()
                    powerpoint =win32com .client .Dispatch ("PowerPoint.Application")
                    powerpoint .Visible =1 
                    try :
                        powerpoint .WindowState =2 
                    except :
                        pass 

                    abs_path =str (file_path .resolve ())
                    ppt =powerpoint .Presentations .Open (abs_path ,ReadOnly =1 ,Untitled =1 ,WithWindow =0 )

                    for slide_idx in range (min (max_slides ,ppt .Slides .Count )):
                        try :
                            slide =ppt .Slides (slide_idx +1 )
                            img_path =PathlibPath (temp_dir )/f"slide_{slide_idx }.png"
                            slide .Export (str (img_path ),"PNG",ScaleWidth =1920 ,ScaleHeight =1080 )

                            _ltime .sleep (0.3 )

                            if img_path .exists ()and img_path .stat ().st_size >0 :
                                with open (img_path ,"rb")as img_file :
                                    img_data =img_file .read ()
                                base64_image =base64 .b64encode (img_data ).decode ("utf-8")

                                html_parts .append (
                                f'<div class="pptx-slide" style="margin: 30px auto; max-width: 95%; text-align: center; border: 1px solid #d1d5db; padding: 20px; border-radius: 12px; background: linear-gradient(to bottom, #ffffff, #f9fafb);"><div class="slide-number" style="font-weight: 600; color: #2e5bff; margin-bottom: 15px;">🎬 슬라이드 {slide_idx +1 } / {total_slides }</div><img src="data:image/png;base64,{base64_image }" style="max-width: 100%; height: auto; border-radius: 4px;"></div>'
                                )
                                try :
                                    os .remove (str (img_path ))
                                except :
                                    pass 
                        except Exception as e :
                            html_parts .append (f'<div style="margin: 30px 0; padding: 20px; border: 2px solid #fee2e2; border-radius: 8px; background: #fef2f2;"><p style="color: #dc2626;">⚠️ 슬라이드 {slide_idx +1 } 오류: {str (e )}</p></div>')

                    if total_slides >max_slides :
                        html_parts .append (f'<p style="color:#6b7280; text-align: center; margin-top: 30px; padding: 20px; background: #f3f4f6; border-radius: 12px;">📚 전체 {total_slides }슬라이드 중 {max_slides }개만 표시됩니다</p>')

                    return "\n".join (html_parts )if html_parts else _fallback_pptx_text (file_path ,max_slides )

                except (ImportError ,Exception ):

                    return _convert_pptx_with_soffice (file_path ,total_slides ,max_slides )

            finally :
                try :
                    if ppt :
                        ppt .Close ()
                    if powerpoint :
                        powerpoint .Quit ()
                except :
                    pass 
                try :
                    import shutil 
                    _ltime .sleep (0.5 )
                    if os .path .exists (temp_dir ):
                        shutil .rmtree (temp_dir ,ignore_errors =True )
                except :
                    pass 


        else :
            return _convert_pptx_with_soffice (file_path ,total_slides ,max_slides )

    except Exception as e :
        return f"<p style='color: #dc2626;'>PowerPoint 파일을 읽을 수 없습니다: {str (e )}</p>"


def _convert_pptx_with_soffice (file_path :Path ,total_slides :int ,max_slides :int )->str :
    """LibreOffice soffice를 사용하여 PPTX 변환 (Linux/Mac)"""
    import tempfile 
    import os 
    import subprocess 
    import glob 
    import time as _ltime 
    import shutil 

    tmp_dir =tempfile .mkdtemp (prefix ="ppt_preview_")
    try :

        try :
            result =subprocess .run (
            ["which","soffice"],
            stdout =subprocess .PIPE ,
            stderr =subprocess .PIPE ,
            timeout =5 
            )
            if result .returncode !=0 :
                return _fallback_pptx_text (file_path ,max_slides )
        except :

            pass 


        cmd =[
        "soffice",
        "--headless",
        "--convert-to","png",
        "--outdir",tmp_dir ,
        str (file_path .resolve ()),
        ]

        subprocess .run (
        cmd ,
        stdout =subprocess .PIPE ,
        stderr =subprocess .PIPE ,
        text =True ,
        timeout =60 ,
        )

        _ltime .sleep (1.0 )


        png_files =sorted (glob .glob (os .path .join (tmp_dir ,"*.png")))
        if not png_files :
            return _fallback_pptx_text (file_path ,max_slides )

        html_parts =[]
        html_parts .append (f'<div style="color: #6b7280; font-size: 12px; margin-bottom: 15px;">총 {total_slides }슬라이드 (이미지 {len (png_files )}장)</div>')

        for idx ,png_path in enumerate (png_files [:max_slides ]):
            try :
                with open (png_path ,"rb")as f :
                    img_data =f .read ()
                b64 =base64 .b64encode (img_data ).decode ("utf-8")
                slide_no =idx +1 

                html_parts .append (
                f'<div class="pptx-slide" style="margin:30px auto;max-width:95%;text-align:center;border:1px solid #d1d5db;padding:20px;border-radius:12px;background:linear-gradient(to bottom,#ffffff,#f9fafb);"><div class="slide-number" style="font-weight:600;color:#2e5bff;margin-bottom:15px;">🎬 슬라이드 {slide_no } / {total_slides }</div><img src="data:image/png;base64,{b64 }" style="max-width:100%;height:auto;border-radius:4px;"></div>'
                )
            except Exception as e :
                html_parts .append (f"<div style='margin:20px 0;padding:15px;border:2px solid #fee2e2;border-radius:8px;background:#fef2f2;color:#dc2626;'>슬라이드 {idx +1 } 오류: {e }</div>")

        if total_slides >max_slides :
            html_parts .append (f"<p style='color:#6b7280;text-align:center;margin-top:20px;padding:15px;background:#f3f4f6;border-radius:8px;'>... 나머지 {total_slides -max_slides }슬라이드 생략</p>")

        return "\n".join (html_parts )if html_parts else "<p>PowerPoint 내용이 없습니다.</p>"

    except subprocess .TimeoutExpired :
        return "<p style='color: #dc2626;'>PowerPoint 변환 시간 초과. LibreOffice 설치 확인: apt-get install libreoffice</p>"
    except FileNotFoundError :
        return "<p style='color: #dc2626;'>LibreOffice가 설치되지 않았습니다. Linux에서는 다음을 실행하세요: sudo apt-get install libreoffice</p>"
    except Exception as e :
        return _fallback_pptx_text (file_path ,max_slides )
    finally :
        try :
            shutil .rmtree (tmp_dir ,ignore_errors =True )
        except :
            pass 


def _fallback_pptx_text (file_path :Path ,max_slides :int )->str :
    try :
        prs =Presentation (file_path )
        total_slides =len (prs .slides )

        html_parts =[]
        html_parts .append (f'<div style="color: #6b7280; font-size: 12px; margin-bottom: 15px;">총 {total_slides }슬라이드 (처음 {max_slides }슬라이드 표시)</div>')

        for slide_idx in range (min (max_slides ,total_slides )):
            try :
                slide =prs .slides [slide_idx ]
                html_parts .append (f'<div class="slide" style="margin: 30px 0; padding: 20px; background: white; border: 2px solid #e5e7eb; border-radius: 8px;">')
                html_parts .append (f'<div class="slide-title" style="font-size: 18px; font-weight: 700; color: #2e5bff; margin-bottom: 20px; padding-bottom: 10px; border-bottom: 2px solid #2e5bff;">🎬 슬라이드 {slide_idx +1 }</div>')

                has_content =False 
                for shape in slide .shapes :
                    try :
                        if hasattr (shape ,"text_frame"):
                            for para in shape .text_frame .paragraphs :
                                if para .text .strip ():
                                    has_content =True 
                                    level =para .level if hasattr (para ,'level')else 0 
                                    font_size =16 if level ==0 else 14 -(level *2 )
                                    font_weight =600 if level ==0 else 400 
                                    margin =f"margin-left: {level *20 }px;"
                                    html_parts .append (f'<p style="margin: 10px 0; line-height: 1.6; {margin }; font-size: {font_size }px; font-weight: {font_weight };">{para .text }</p>')

                        if hasattr (shape ,"table"):
                            has_content =True 
                            table =shape .table 
                            html_parts .append ('<table style="border-collapse: collapse; width: 100%; margin: 15px 0; border: 1px solid #d1d5db;">')
                            for row_idx ,row in enumerate (table .rows ):
                                html_parts .append ('<tr>')
                                for cell in row .cells :
                                    bg_color ="background-color: #f3f4f6; font-weight: 600;"if row_idx ==0 else ""
                                    html_parts .append (f'<td style="border: 1px solid #d1d5db; padding: 10px; {bg_color }">')
                                    for para in cell .text_frame .paragraphs :
                                        if para .text .strip ():
                                            html_parts .append (f'<p style="margin: 0;">{para .text }</p>')
                                    html_parts .append ('</td>')
                                html_parts .append ('</tr>')
                            html_parts .append ('</table>')
                    except :
                        pass 

                if not has_content :
                    html_parts .append ('<p style="color: #9ca3af; font-style: italic;">이 슬라이드에는 표시할 텍스트 콘텐츠가 없습니다.</p>')

                html_parts .append ('</div>')
            except Exception as e :
                html_parts .append (f'<div class="slide"><p style="color: #dc2626;">슬라이드 {slide_idx +1 } 읽기 오류: {str (e )}</p></div>')

        if total_slides >max_slides :
            html_parts .append (f'<p style="color:#6b7280; text-align: center; margin-top: 20px; padding: 15px; background: #f3f4f6; border-radius: 8px;">... 나머지 {total_slides -max_slides }슬라이드 생략</p>')

        return '\n'.join (html_parts )if html_parts else "<p>PowerPoint 내용이 없습니다.</p>"
    except Exception as e :
        return f"<p>PowerPoint 파일을 읽을 수 없습니다: {str (e )}</p>"


def preview_image (file_path :Path )->str :
    """이미지 파일 미리보기"""
    try :

        if not file_path .exists ():
            return f'<p style="color: #dc2626;">이미지 파일을 찾을 수 없습니다: {file_path }</p>'


        with open (file_path ,'rb')as f :
            img_data =f .read ()

        base64_image =base64 .b64encode (img_data ).decode ('utf-8')


        ext =file_path .suffix .lower ()
        mime_types ={
        '.jpg':'image/jpeg',
        '.jpeg':'image/jpeg',
        '.png':'image/png',
        '.gif':'image/gif',
        '.webp':'image/webp',
        '.svg':'image/svg+xml',
        '.bmp':'image/bmp',
        }
        mime_type =mime_types .get (ext ,'image/jpeg')

        return f'<img src="data:{mime_type };base64,{base64_image }" alt="Image" style="max-width:100%;height:auto;border-radius:8px;">'
    except Exception as e :
        return f'<p style="color: #dc2626;">이미지를 읽을 수 없습니다: {str (e )}</p>'


def preview_text (file_path :Path )->str :
    try :
        content =""
        for encoding in ['utf-8','cp949','euc-kr','latin-1']:
            try :
                with open (file_path ,'r',encoding =encoding )as f :
                    content =f .read (10000 )
                break 
            except (UnicodeDecodeError ,LookupError ):
                continue 

        if not content :
            content ="[파일을 텍스트로 읽을 수 없습니다]"

        content =content .replace ('&','&amp;').replace ('<','&lt;').replace ('>','&gt;')
        return f'<div style="background:#f9fafb;padding:15px;border-radius:8px;border:1px solid #e5e7eb;white-space:pre-wrap;word-break:break-word;font-family:\'Courier New\',monospace;font-size:12px;max-height:600px;overflow-y:auto;color:#1f2937;">{content }</div>'
    except Exception as e :
        return f'<p>파일을 읽을 수 없습니다: {str (e )}</p>'


def get_file_preview (file_path :Path )->str :
    ext =file_path .suffix .lower ()

    if ext in ['.xlsx','.xls']:
        return preview_excel (file_path )
    elif ext in ['.docx','.doc']:
        return preview_word (file_path )
    elif ext =='.pdf':
        return preview_pdf (file_path )
    elif ext in ['.pptx','.ppt']:
        return preview_pptx (file_path )
    elif ext in ['.jpg','.jpeg','.png','.gif','.webp','.svg','.bmp']:
        return preview_image (file_path )
    else :
        return preview_text (file_path )