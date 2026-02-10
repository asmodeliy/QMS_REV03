"""
RAG Document Indexer for QMS

This module indexes all documents in the QMS project for use with the GPT4All assistant.
It creates a knowledge base that can be retrieved and used in prompts.
"""

import os
import sqlite3
from pathlib import Path
from typing import List, Dict, Any, Optional
import hashlib
import json
from datetime import datetime


class RAGIndexer:
    """Index documents from the QMS project for RAG retrieval"""
    
    def __init__(self, db_path: str = "rag_knowledge_base.db", base_dir: Optional[Path] = None):
        """Initialize the RAG indexer
        
        Args:
            db_path: Path to the SQLite database for storing indexed documents
            base_dir: Base directory of the QMS project
        """
        self.db_path = db_path
        self.base_dir = base_dir or Path(__file__).resolve().parent.parent.parent
        self._init_db()
    
    def _init_db(self):
        """Initialize or verify the SQLite database schema"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        # Create tables
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS documents (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                file_path TEXT UNIQUE NOT NULL,
                file_name TEXT NOT NULL,
                file_type TEXT,
                content_hash TEXT NOT NULL,
                content TEXT NOT NULL,
                summary TEXT,
                keywords TEXT,
                module_name TEXT,
                indexed_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        """)
        
        cursor.execute("""
            CREATE TABLE IF NOT EXISTS chunks (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                document_id INTEGER NOT NULL,
                chunk_index INTEGER NOT NULL,
                chunk_text TEXT NOT NULL,
                chunk_keywords TEXT,
                FOREIGN KEY (document_id) REFERENCES documents(id) ON DELETE CASCADE,
                UNIQUE(document_id, chunk_index)
            )
        """)
        
        cursor.execute("""
            CREATE INDEX IF NOT EXISTS idx_file_path ON documents(file_path)
        """)
        
        cursor.execute("""
            CREATE INDEX IF NOT EXISTS idx_module_name ON documents(module_name)
        """)
        
        conn.commit()
        conn.close()
    
    def _get_file_hash(self, file_path: Path) -> str:
        """Get SHA256 hash of file content"""
        try:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            return hashlib.sha256(content.encode()).hexdigest()
        except Exception:
            return ""
    
    def _extract_summary(self, content: str, max_length: int = 200) -> str:
        """Extract a summary from document content"""
        lines = content.split('\n')
        # Try to get first docstring or comment block
        for i, line in enumerate(lines[:20]):
            if '"""' in line or "'''" in line or '#' in line:
                summary_lines = []
                for j in range(i, min(i + 5, len(lines))):
                    l = lines[j].strip()
                    if l:
                        summary_lines.append(l)
                summary = ' '.join(summary_lines)
                return summary[:max_length]
        
        # Fall back to first non-empty line
        for line in lines:
            stripped = line.strip()
            if stripped and not stripped.startswith('#'):
                return stripped[:max_length]
        return ""
    
    def _extract_keywords(self, content: str, file_name: str) -> List[str]:
        """Extract keywords from document"""
        keywords = set()
        
        # Add file name parts
        name_parts = file_name.replace('.py', '').replace('_', ' ').split()
        keywords.update(name_parts)
        
        # Extract class and function names
        lines = content.split('\n')
        for line in lines:
            line = line.strip()
            if line.startswith('class '):
                class_name = line.split('class ')[1].split('(')[0].split(':')[0].strip()
                keywords.add(class_name)
            elif line.startswith('def '):
                func_name = line.split('def ')[1].split('(')[0].strip()
                keywords.add(func_name)
            elif line.startswith('async def '):
                func_name = line.split('async def ')[1].split('(')[0].strip()
                keywords.add(func_name)
        
        return list(keywords)[:20]
    
    def _chunk_content(self, content: str, chunk_size: int = 1000) -> List[str]:
        """Split content into chunks"""
        chunks = []
        for i in range(0, len(content), chunk_size):
            chunk = content[i:i + chunk_size].strip()
            if chunk:
                chunks.append(chunk)
        return chunks if chunks else [content]
    
    def index_file(self, file_path: Path, module_name: Optional[str] = None) -> bool:
        """Index a single file (supports .py, .html, .js, .css, .md, .docx, .pdf, .pptx)
        
        Args:
            file_path: Path to the file to index
            module_name: Module name/category for the document
        
        Returns:
            True if successfully indexed, False otherwise
        """
        try:
            if not file_path.exists():
                return False
            
            # Read file content based on file type
            file_ext = file_path.suffix.lower()
            content = ""
            
            if file_ext in ['.py', '.html', '.js', '.css', '.md', '.txt', '.sql']:
                # Text files
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            
            elif file_ext == '.docx':
                # Word documents
                try:
                    from docx import Document
                    doc = Document(file_path)
                    content = "\n".join([para.text for para in doc.paragraphs])
                except ImportError:
                    print(f"Warning: python-docx not installed. Skipping {file_path}")
                    return False
                except Exception as e:
                    print(f"Warning: Could not read DOCX {file_path}: {e}")
                    return False
            
            elif file_ext == '.pdf':
                # PDF files
                try:
                    from pypdf import PdfReader
                except ImportError:
                    from PyPDF2 import PdfReader
                
                try:
                    reader = PdfReader(file_path)
                    content = "\n".join([page.extract_text() for page in reader.pages])
                except Exception as e:
                    print(f"Warning: Could not read PDF {file_path}: {e}")
                    return False
            
            elif file_ext == '.pptx':
                # PowerPoint presentations
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    slides_text = []
                    for slide in prs.slides:
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text:
                                slide_text.append(shape.text)
                        if slide_text:
                            slides_text.append("\n".join(slide_text))
                    content = "\n\n---\n\n".join(slides_text)
                except ImportError:
                    print(f"Warning: python-pptx not installed. Skipping {file_path}")
                    return False
                except Exception as e:
                    print(f"Warning: Could not read PPTX {file_path}: {e}")
                    return False
            
            else:
                # Unsupported file type
                return False
            
            if not content or len(content.strip()) < 10:
                return False
            
            file_hash = self._get_file_hash(file_path)
            relative_path = file_path.relative_to(self.base_dir)
            
            # Extract metadata
            summary = self._extract_summary(content)
            keywords = self._extract_keywords(content, file_path.name)
            chunks = self._chunk_content(content)
            
            conn = sqlite3.connect(self.db_path)
            cursor = conn.cursor()
            
            try:
                # Check if document already indexed and hasn't changed
                cursor.execute(
                    "SELECT id, content_hash FROM documents WHERE file_path = ?",
                    (str(relative_path),)
                )
                result = cursor.fetchone()
                
                if result and result[1] == file_hash:
                    # Already indexed and content hasn't changed
                    return True
                
                # Delete old document if exists
                if result:
                    cursor.execute("DELETE FROM documents WHERE id = ?", (result[0],))
                
                # Insert new document
                cursor.execute("""
                    INSERT INTO documents 
                    (file_path, file_name, file_type, content_hash, content, summary, keywords, module_name)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                """, (
                    str(relative_path),
                    file_path.name,
                    file_path.suffix,
                    file_hash,
                    content[:50000],  # Limit content size
                    summary,
                    json.dumps(keywords),
                    module_name or file_ext[1:].upper()
                ))
                
                doc_id = cursor.lastrowid
                
                # Insert chunks
                for idx, chunk in enumerate(chunks):
                    chunk_keywords = self._extract_keywords(chunk, file_path.name)
                    cursor.execute("""
                        INSERT INTO chunks (document_id, chunk_index, chunk_text, chunk_keywords)
                        VALUES (?, ?, ?, ?)
                    """, (doc_id, idx, chunk, json.dumps(chunk_keywords)))
                
                conn.commit()
                return True
                
            finally:
                conn.close()
                
        except Exception as e:
            print(f"Error indexing {file_path}: {e}")
            return False
    
    def index_directory(self, dir_path: Path, pattern: str = "*.py", recursive: bool = True) -> int:
        """Index all files in a directory
        
        Args:
            dir_path: Directory to index
            pattern: File pattern to match (e.g., "*.py")
            recursive: If True, recursively search subdirectories
        
        Returns:
            Number of files successfully indexed
        """
        if not dir_path.exists():
            return 0
        
        count = 0
        skip_dirs = {'__pycache__', '.git', 'node_modules', '.venv', 'venv', 'build', 'dist'}
        
        if recursive:
            for root, dirs, files in os.walk(dir_path):
                # Skip certain directories
                dirs[:] = [d for d in dirs if d not in skip_dirs]
                
                # Get module name from path
                rel_root = Path(root).relative_to(dir_path)
                module_name = str(rel_root).replace(os.sep, '.') if str(rel_root) != '.' else ''
                
                for file in files:
                    if Path(file).match(pattern):
                        file_path = Path(root) / file
                        if self.index_file(file_path, module_name):
                            count += 1
        else:
            module_name = dir_path.name
            for file_path in dir_path.glob(pattern):
                if file_path.is_file():
                    if self.index_file(file_path, module_name):
                        count += 1
        
        return count
    
    def index_project(self) -> Dict[str, int]:
        """Index the entire QMS project including spec-center documents
        
        Returns:
            Dictionary with count of indexed files by type
        """
        counts = {
            'python': 0,
            'html': 0,
            'javascript': 0,
            'css': 0,
            'sql': 0,
            'markdown': 0,
            'docx': 0,
            'pdf': 0,
            'pptx': 0,
            'total': 0
        }
        
        # Index Python files
        counts['python'] += self.index_directory(self.base_dir, "*.py")
        
        # Index template files
        templates_dir = self.base_dir / "templates"
        counts['html'] += self.index_directory(templates_dir, "*.html")
        
        # Index static files (JS, CSS)
        static_dir = self.base_dir / "static"
        counts['javascript'] += self.index_directory(static_dir, "*.js")
        counts['css'] += self.index_directory(static_dir, "*.css")
        
        # Index markdown documentation
        docs_dir = self.base_dir / "docs"
        counts['markdown'] += self.index_directory(docs_dir, "*.md")
        
        # Index spec-center documents
        spec_center_dir = self.base_dir / "modules" / "spec_center" / "uploads"
        if spec_center_dir.exists():
            print(f"[RAG] Indexing spec-center documents from {spec_center_dir}")
            counts['docx'] += self.index_directory(spec_center_dir, "*.docx")
            counts['pdf'] += self.index_directory(spec_center_dir, "*.pdf")
            counts['pptx'] += self.index_directory(spec_center_dir, "*.pptx")
        
        # Index uploads directory 
        uploads_dir = self.base_dir / "uploads"
        if uploads_dir.exists():
            print(f"[RAG] Indexing uploads from {uploads_dir}")
            counts['docx'] += self.index_directory(uploads_dir, "*.docx")
            counts['pdf'] += self.index_directory(uploads_dir, "*.pdf")
            counts['pptx'] += self.index_directory(uploads_dir, "*.pptx")
        
        counts['total'] = sum(v for k, v in counts.items() if k != 'total')
        return counts
    
    def search(self, query: str, limit: int = 5) -> List[Dict[str, Any]]:
        """Search for relevant documents
        
        Args:
            query: Search query
            limit: Maximum number of results to return
        
        Returns:
            List of matching documents with relevance scores
        """
        query_terms = set(query.lower().split())
        
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        try:
            # Get all documents and calculate relevance
            cursor.execute("SELECT id, file_path, file_name, content, summary, keywords FROM documents")
            all_docs = cursor.fetchall()
            
            results = []
            for doc_id, file_path, file_name, content, summary, keywords_json in all_docs:
                try:
                    keywords = json.loads(keywords_json) if keywords_json else []
                except:
                    keywords = []
                
                # Calculate relevance score
                score = 0
                
                # File name match
                file_terms = set(file_name.lower().replace('_', ' ').replace('.', ' ').split())
                score += len(query_terms & file_terms) * 3
                
                # Keywords match
                keyword_terms = set(k.lower() for k in keywords)
                score += len(query_terms & keyword_terms) * 2
                
                # Content match (simple substring search)
                content_lower = content.lower()
                for term in query_terms:
                    score += content_lower.count(term)
                
                # Summary match
                if summary:
                    summary_lower = summary.lower()
                    score += len(query_terms & set(summary_lower.split())) * 2
                
                if score > 0:
                    results.append({
                        'id': doc_id,
                        'file_path': file_path,
                        'file_name': file_name,
                        'summary': summary,
                        'content': content[:500],  # Preview
                        'keywords': keywords,
                        'score': score
                    })
            
            # Sort by score and return top results
            results.sort(key=lambda x: x['score'], reverse=True)
            return results[:limit]
            
        finally:
            conn.close()
    
    def get_context(self, query: str, max_tokens: int = 2000) -> str:
        """Get context string for LLM prompt
        
        Args:
            query: Query to find relevant documents
            max_tokens: Maximum tokens to include in context
        
        Returns:
            Formatted context string
        """
        results = self.search(query, limit=10)
        
        context_parts = []
        total_tokens = 0
        
        for result in results:
            # Estimate tokens (rough: 1 token ≈ 4 chars)
            doc_size = len(result['file_path']) + len(result['summary']) + len(result['content'])
            estimated_tokens = doc_size // 4
            
            if total_tokens + estimated_tokens > max_tokens:
                break
            
            context_parts.append(f"""
Document: {result['file_path']}
Summary: {result['summary']}
Keywords: {', '.join(result['keywords'][:5])}
Content Preview:
{result['content']}
---""")
            total_tokens += estimated_tokens
        
        if not context_parts:
            return ""
        
        return f"Retrieved Knowledge Base Context:\n{''.join(context_parts)}"
    
    def get_stats(self) -> Dict[str, Any]:
        """Get indexing statistics"""
        conn = sqlite3.connect(self.db_path)
        cursor = conn.cursor()
        
        try:
            cursor.execute("SELECT COUNT(*) FROM documents")
            total_docs = cursor.fetchone()[0]
            
            cursor.execute("SELECT COUNT(*) FROM chunks")
            total_chunks = cursor.fetchone()[0]
            
            cursor.execute("SELECT module_name, COUNT(*) FROM documents GROUP BY module_name")
            module_stats = dict(cursor.fetchall())
            
            return {
                'total_documents': total_docs,
                'total_chunks': total_chunks,
                'modules': module_stats,
                'db_path': str(self.db_path)
            }
        finally:
            conn.close()


if __name__ == "__main__":
    # Example usage
    indexer = RAGIndexer()
    print("Indexing QMS project...")
    counts = indexer.index_project()
    print(f"Indexed {counts}")
    
    print("\nIndexing statistics:")
    stats = indexer.get_stats()
    print(f"Total documents: {stats['total_documents']}")
    print(f"Total chunks: {stats['total_chunks']}")
    print(f"Modules: {stats['modules']}")
    
    # Example search
    print("\nExample search for 'project':")
    results = indexer.search("project", limit=3)
    for result in results:
        print(f"  - {result['file_path']} (score: {result['score']})")
